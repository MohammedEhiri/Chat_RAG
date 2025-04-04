# --- Standard Library Imports ---
import sys
import os
from pathlib import Path
from typing import List, Dict, Optional, Generator, Tuple, Any, AsyncGenerator, Union
import asyncio
import logging
import json
import time
import hashlib
from concurrent.futures import ThreadPoolExecutor
import tempfile
from io import BytesIO
import re # <-- Added import
import torch
from langchain_core.runnables import Runnable
import argparse
import shutil
from operator import itemgetter

# Data processing imports
from pypdf import PdfReader
from pptx import Presentation
from docx import Document as WordDocument
import nltk
from PIL import Image
# Make OCR optional
try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logging.warning("pytesseract not found. Install it and Tesseract OCR for image text extraction.")

# Vector DB and embeddings
from langchain_chroma import Chroma
# from langchain.vectorstores import FAISS # Not used currently
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.embeddings import FastEmbedEmbeddings, CacheBackedEmbeddings # FastEmbed not used
from langchain.storage import LocalFileStore # <-- Use file store for embedding cache persistence
from langchain.embeddings import CacheBackedEmbeddings as LangchainCacheBackedEmbeddings # <-- Correct import path for CacheBackedEmbeddings used with store

# Text processing
from langchain_text_splitters import (
    RecursiveCharacterTextSplitter,
    MarkdownHeaderTextSplitter,
    CharacterTextSplitter # Not explicitly used, could remove
)

# LLM and chains
from langchain_ollama import OllamaLLM
from langchain_core.output_parsers import StrOutputParser, JsonOutputParser # JsonOutputParser not used currently
from langchain_core.prompts import ChatPromptTemplate, PromptTemplate
from langchain_core.runnables import RunnablePassthrough, RunnableLambda, RunnableParallel, RunnableConfig, ensure_config
from langchain.retrievers import BM25Retriever, EnsembleRetriever
from langchain.retrievers.document_compressors import LLMChainExtractor
# from langchain.chains.combine_documents import create_stuff_documents_chain # Not used directly with LCEL approach
from langchain.retrievers import ContextualCompressionRetriever # <-- Correct import for ContextualCompressionRetriever
# from langchain_core.callbacks import AsyncCallbackManager # Not used
from langchain_community.cross_encoders import HuggingFaceCrossEncoder

# Reranker
from langchain.retrievers.document_compressors import CrossEncoderReranker
from sentence_transformers import CrossEncoder # Keep for model loading

# UI
import gradio as gr
from gradio.themes.utils import colors, sizes # Not used currently
from gradio.themes import Base as ThemeBase # Not used currently

# Initialize NLTK for better text processing
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)

try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords', quiet=True)

from nltk.corpus import stopwords

# --- Configuration ---
EMBEDDING_MODEL = "sentence-transformers/all-mpnet-base-v2" # Faster than BGE, good starting point
LLM_MODEL = "gemma3:1b" # As specified
DB_DIR = Path("vectore/enhanced_rag_db_v2") # New dir for clarity
DATA_DIR = Path("documents") # As specified
CACHE_METADATA_FILE = DB_DIR / "rag_cache_metadata.json"
CHUNK_CACHE_FILE = DB_DIR / "chunk_cache.json" # <-- File to store chunks for BM25 loading
RESULT_CACHE_FILE = DB_DIR / "query_cache.json"
EMBEDDING_CACHE_DIR = DB_DIR / "embedding_cache" # <-- Directory for embedding cache persistence

RETRIEVER_DENSE_K = 15 # K for dense retriever before ensemble
RETRIEVER_SPARSE_K = 15 # K for sparse retriever before ensemble
ENSEMBLE_RANK_K = 15 # How many results after ensemble rank fusion (RRF)
RERANKER_TOP_N = 8 # How many docs after cross-encoder reranking
RERANKER_MODEL = "BAAI/bge-reranker-large" # Powerful reranker

CHUNK_SIZE = 800
CHUNK_OVERLAP = 150

# --- Logging Setup ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger("EnhancedRAG")

# --- Explicit GPU Check ---
GPU_DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
logger.info(f"--- Device Check ---")
logger.info(f"Using device: {GPU_DEVICE}")
if GPU_DEVICE == "cuda":
    try:
        logger.info(f"PyTorch version: {torch.__version__}")
        logger.info(f"CUDA available: {torch.cuda.is_available()}")
        logger.info(f"CUDA version detected by PyTorch: {torch.version.cuda}")
        logger.info(f"Number of GPUs: {torch.cuda.device_count()}")
        for i in range(torch.cuda.device_count()):
            logger.info(f"GPU {i}: {torch.cuda.get_device_name(i)}")
    except Exception as e:
        logger.error(f"Error during GPU check: {e}")
else:
    logger.warning("CUDA not available, running on CPU.")
logger.info(f"--------------------")

# --- Helper Classes (ResponseCache, DocumentProcessor, EnhancedDocumentSplitter, QueryProcessor) ---
# Keep the implementations provided by the user for these classes, they look good.
# (Minor tweak: Add import re to QueryProcessor if not already present)
# ... (Paste User's Class Definitions Here) ...
class ResponseCache:
    """Cache for query responses to improve performance"""
    def __init__(self, cache_file: Path):
        self.cache_file = cache_file
        self.cache = self._load_cache()
        self.max_cache_size = 1000  # Maximum number of entries

    def _load_cache(self) -> Dict:
        if self.cache_file.exists():
            try:
                with open(self.cache_file, 'r') as f:
                    return json.load(f)
            except Exception as e:
                logger.warning(f"Failed to load response cache: {e}")
                return {}
        return {}

    def save_cache(self):
        """Save cache to disk"""
        # Ensure directory exists
        self.cache_file.parent.mkdir(parents=True, exist_ok=True)
        try:
            # Add timestamp for potential LRU logic later if needed
            for k in self.cache:
                if "timestamp" not in self.cache[k]:
                     self.cache[k]["timestamp"] = time.time()

            with open(self.cache_file, 'w') as f:
                json.dump(self.cache, f, indent=2) # Add indent for readability
        except Exception as e:
            logger.error(f"Failed to save response cache: {e}")

    def get(self, query: str, history_hash: str = "") -> Optional[Dict]:
        """Get cached response for query and history"""
        key = self._create_key(query, history_hash)
        result = self.cache.get(key)
        if result:
             logger.info(f"Response cache HIT for query: '{query[:50]}...'")
             # Update timestamp on access? Optional for LRU.
             # result['timestamp'] = time.time()
        else:
             logger.info(f"Response cache MISS for query: '{query[:50]}...'")
        return result

    def set(self, query: str, history_hash: str, response: Dict):
        """Cache response for query and history"""
        if not isinstance(response, dict):
             logger.error("Attempted to cache non-dict response. Skipping.")
             return
        key = self._create_key(query, history_hash)
        response_to_cache = response.copy() # Avoid modifying original dict
        response_to_cache["timestamp"] = time.time() # Add timestamp for LRU
        self.cache[key] = response_to_cache

        # Trim cache if needed using simple FIFO based on insertion order (approx LRU)
        if len(self.cache) > self.max_cache_size:
            logger.info(f"Cache size ({len(self.cache)}) exceeds max ({self.max_cache_size}). Trimming...")
            # Get keys sorted by timestamp (oldest first)
            keys_to_remove = sorted(self.cache.keys(), key=lambda k: self.cache[k].get("timestamp", 0))
            num_to_remove = len(self.cache) - self.max_cache_size + 50 # Remove a bit extra
            for i in range(min(num_to_remove, len(keys_to_remove))):
                removed_key = keys_to_remove[i]
                self.cache.pop(removed_key, None)
                logger.debug(f"Removed oldest cache entry: {removed_key}")

        self.save_cache()

    def _create_key(self, query: str, history_hash: str) -> str:
        """Create a cache key from query and history hash"""
        query_normalized = query.lower().strip()
        combined = f"{query_normalized}||{history_hash}"
        return hashlib.md5(combined.encode()).hexdigest()

    def create_history_hash(self, chat_history: List[Tuple[str, str]]) -> str:
        """Create a hash from chat history"""
        if not chat_history:
            return ""
        # Use last N exchanges to capture relevant recent context
        history_context = chat_history[-3:] # Consider last 3 exchanges
        try:
            history_str = json.dumps(history_context, sort_keys=True) # Ensure consistent order
        except TypeError:
            logger.warning("Could not serialize chat history for hashing. Using empty hash.")
            return ""
        return hashlib.md5(history_str.encode()).hexdigest()


class DocumentProcessor:
    """Handles processing of various document types"""

    def __init__(self, data_dir: Path):
        self.data_dir = data_dir
        self.ocr_available = OCR_AVAILABLE # Use global check

    def process_pdf(self, file_path: Path) -> List[Document]:
        """Process PDF files with improved text and image extraction"""
        documents = []
        try:
            logger.info(f"Processing PDF: {file_path.name}")
            pdf = PdfReader(file_path)
            num_pages = len(pdf.pages)
            if num_pages == 0:
                 logger.warning(f"Skipping PDF with 0 pages: {file_path.name}")
                 return []
            pdf_metadata_base = {
                "source": str(file_path.relative_to(self.data_dir, walk_up=True)) if self.data_dir in file_path.parents else file_path.name, # Handle relative path better
                "file_name": file_path.name,
                "file_type": ".pdf",
                "total_pages": num_pages
            }

            for page_idx, page in enumerate(pdf.pages):
                page_num = page_idx + 1
                page_metadata = {**pdf_metadata_base, "page_number": page_num}
                text = ""
                try:
                    text = page.extract_text() or "" # Ensure text is not None
                except Exception as text_extract_e:
                     logger.warning(f"Could not extract text from page {page_num} of {file_path.name}: {text_extract_e}")

                image_text = ""
                if self.ocr_available:
                    try:
                        # Using page.images is efficient with pypdf
                        count = 0
                        for img_obj in page.images:
                            count += 1
                            try:
                                # Convert Pillow Image from pypdf Image object
                                img = img_obj.to_pil()
                                ocr_text = pytesseract.image_to_string(img)
                                if ocr_text.strip():
                                    image_text += f"\n[Image {count} on Page {page_num} Content]:\n{ocr_text.strip()}\n"
                            except Exception as ocr_e:
                                logger.debug(f"OCR failed for image {count} on page {page_num} in {file_path.name}: {ocr_e}")
                    except Exception as img_extract_e:
                         logger.warning(f"Could not extract images from page {page_num} of {file_path.name}: {img_extract_e}")

                combined_text = (text.strip() + "\n" + image_text.strip()).strip()

                if not combined_text:
                    logger.debug(f"  Skipping page {page_num} in {file_path.name} (no text found)")
                    continue

                # Add context marker
                contextual_text = f"[PDF Page {page_num}/{num_pages}]\n{combined_text}"
                doc = Document(page_content=contextual_text, metadata=page_metadata)
                documents.append(doc)
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}", exc_info=True)

        return documents

    def process_pptx(self, file_path: Path) -> List[Document]:
        """Process PowerPoint files with improved text extraction"""
        documents = []
        try:
            logger.info(f"Processing presentation: {file_path.name}")
            presentation = Presentation(file_path)
            num_slides = len(presentation.slides)
            if num_slides == 0:
                 logger.warning(f"Skipping PPTX with 0 slides: {file_path.name}")
                 return []
            pptx_metadata_base = {
                "source": str(file_path.relative_to(self.data_dir, walk_up=True)) if self.data_dir in file_path.parents else file_path.name,
                "file_name": file_path.name,
                "file_type": ".pptx",
                "total_slides": num_slides
            }

            for slide_idx, slide in enumerate(presentation.slides):
                slide_num = slide_idx + 1
                slide_metadata = {**pptx_metadata_base, "slide_number": slide_num}

                all_text_content, slide_title = self._extract_all_text_from_slide(slide)

                if slide_title:
                     slide_metadata["slide_title"] = slide_title

                if not all_text_content or not all_text_content.strip():
                    logger.debug(f"  Skipping empty slide {slide_num} in {file_path.name}")
                    continue

                # Add context marker
                title_marker = f" - Title: {slide_title}" if slide_title else ""
                contextual_text = f"[Slide {slide_num}/{num_slides}{title_marker}]\n{all_text_content}"

                doc = Document(page_content=contextual_text, metadata=slide_metadata)
                documents.append(doc)
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}", exc_info=True)

        return documents

    def _extract_all_text_from_slide(self, slide) -> Tuple[str, Optional[str]]:
        """Extracts text from shapes, tables, notes and identifies title."""
        text_parts = []
        slide_title = None

        # Extract notes first
        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_text = f"[Slide Notes]:\n{notes}\n"
        except Exception as e:
            logger.warning(f"Could not extract notes from slide: {e}")

        # Process shapes
        body_texts = []
        table_texts = []
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        # Simple title heuristic (often placeholder shapes)
                        if shape.shape_type == 1 and not slide_title: # Check for specific title placeholder type MSO_SHAPE_TYPE.PLACEHOLDER (14) -> TITLE (1)
                             # Refine this heuristic if needed based on common templates
                             if "title" in shape.name.lower() or shape.text_frame.text.isupper(): # Simple heuristics
                                 slide_title = text
                                 continue # Don't add title to body text

                        body_texts.append(text)
                elif shape.has_table:
                    table_str = self._extract_table_text(shape.table)
                    if table_str:
                        table_texts.append(table_str)
        except Exception as e:
            logger.warning(f"Error extracting text from slide shapes: {e}")

        # Assemble content
        if slide_title:
            text_parts.append(f"[Title]: {slide_title}")
        if body_texts:
            text_parts.append("[Content]:\n" + "\n".join(body_texts))
        if table_texts:
            text_parts.append("\n".join(table_texts)) # Table extraction adds "Table:" prefix
        if notes_text:
            text_parts.append(notes_text)

        return "\n\n".join(text_parts).strip(), slide_title

    def _extract_table_text(self, table) -> str:
         """Helper to extract text from a table shape."""
         table_rows = []
         header_row = []
         try:
              # Try to identify header row
              if table.rows:
                   first_row = table.rows[0]
                   if all(cell.text_frame and cell.text_frame.text.strip() for cell in first_row.cells):
                        header_row = [cell.text_frame.text.strip() for cell in first_row.cells]

              # Process all rows
              for row_idx, row in enumerate(table.rows):
                   if row_idx == 0 and header_row: continue # Skip header if identified

                   cell_texts = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame and cell.text_frame.text.strip()]
                   if cell_texts:
                        table_rows.append(" | ".join(cell_texts))

              if header_row or table_rows:
                   table_text = "[Table]:\n"
                   if header_row:
                        table_text += f"  Headers: {' | '.join(header_row)}\n"
                   if table_rows:
                        table_text += "  Data:\n" + "\n".join([f"    - {r}" for r in table_rows])
                   return table_text
         except Exception as e:
              logger.warning(f"Could not process table content: {e}")
         return ""


    def process_docx(self, file_path: Path) -> List[Document]:
        """Process Word documents with structure preservation"""
        documents = []
        try:
            logger.info(f"Processing Word document: {file_path.name}")
            doc = WordDocument(file_path)
            doc_title = ""
            try:
                 if doc.core_properties.title:
                      doc_title = doc.core_properties.title
            except AttributeError: pass # Handle cases where properties might be missing

            word_metadata = {
                "source": str(file_path.relative_to(self.data_dir, walk_up=True)) if self.data_dir in file_path.parents else file_path.name,
                "file_name": file_path.name,
                "file_type": ".docx",
                "title": doc_title if doc_title else None
            }

            full_text_parts = []
            if doc_title:
                 full_text_parts.append(f"[Document Title]: {doc_title}")

            current_heading = ""
            for element in doc.element.body:
                 if element.tag.endswith('p'):
                      para = element # Simplified access, might need refinement based on python-docx version
                      style_name = para.style.name if hasattr(para, 'style') and hasattr(para.style, 'name') else ''
                      text = para.text.strip()
                      if not text: continue

                      is_heading = 'heading' in style_name.lower()
                      if is_heading:
                           # Simple heading detection based on style name
                           heading_level = 1
                           try:
                                level_str = style_name.split()[-1]
                                if level_str.isdigit():
                                     heading_level = int(level_str)
                           except: pass
                           current_heading = f"{'#' * heading_level} {text}"
                           full_text_parts.append(current_heading)
                      else:
                           full_text_parts.append(text)

                 elif element.tag.endswith('tbl'):
                      table = element # Simplified access
                      table_str = self._extract_docx_table_text(table, doc) # Pass doc to resolve relationships if needed
                      if table_str:
                           full_text_parts.append(table_str)

            combined_text = "\n\n".join(full_text_parts).strip()
            if not combined_text:
                logger.debug(f"  Skipping empty document: {file_path.name}")
                return documents

            doc_entry = Document(page_content=combined_text, metadata=word_metadata)
            documents.append(doc_entry)
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}", exc_info=True)

        return documents

    def _extract_docx_table_text(self, table_element, doc_part) -> str:
         """Helper to extract text from a docx table element."""
         # This requires interacting with the low-level XML or using a more specific
         # python-docx approach if available for the element structure.
         # For simplicity, let's fall back to iterating paragraphs within the table context
         # NOTE: This fallback is basic and might miss complex table structures.
         # A better implementation would parse the w:tbl XML directly.
         table_text_content = []
         try:
             # Find the Table object corresponding to the element (this is complex)
             # Find the Table object corresponding to the element
             # This is tricky - python-docx doesn't directly map body elements back to Table objects easily.
             # We'll attempt a basic text extraction from paragraphs within the element's XML.
             for para_element in table_element.xpath('.//w:p', namespaces=table_element.nsmap):
                  # Reconstruct paragraph text (basic)
                  para_text = "".join(node.text for node in para_element.xpath('.//w:t', namespaces=para_element.nsmap) if node.text)
                  if para_text.strip():
                       table_text_content.append(para_text.strip())

             if table_text_content:
                  # Basic formatting, doesn't preserve rows/columns well
                  return "[Table]:\n" + "\n".join(f"  - {line}" for line in table_text_content)

         except Exception as e:
              logger.warning(f"Could not process complex DOCX table structure: {e}")
         return ""


class EnhancedDocumentSplitter:
    """Advanced document splitter with semantic and structural awareness"""

    def __init__(self, chunk_size: int = CHUNK_SIZE, chunk_overlap: int = CHUNK_OVERLAP):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

        self.recursive_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", ", ", " ", ""], # Common separators
            length_function=len,
        )
        self.header_splitter = MarkdownHeaderTextSplitter(
            headers_to_split_on=[
                ("#", "header_1"), # Assumes markdown H1
                ("##", "header_2"),# Assumes markdown H2
                ("###", "header_3"),# Assumes markdown H3
            ],
            return_each_line=False # Keep header content together
        )

    def split_documents(self, documents: List[Document]) -> List[Document]:
        """Split documents based on their type and content"""
        all_chunks = []
        doc_index = 0 # Keep track of original document index

        for doc in documents:
            file_type = doc.metadata.get("file_type", "")
            content = doc.page_content
            original_metadata = doc.metadata.copy()
            original_metadata["original_doc_index"] = doc_index

            # Heuristic to check for markdown-like headers
            header_pattern = r'^\s*#{1,3}\s+'
            has_headers = bool(re.search(header_pattern, content, re.MULTILINE))

            split_docs = []
            if file_type == ".docx" and has_headers: # Apply header splitter primarily to DOCX where we added #
                logger.debug(f"Applying Header Splitter to: {original_metadata.get('file_name')}")
                try:
                    header_splits = self.header_splitter.split_text(content)
                    for split in header_splits:
                         merged_metadata = original_metadata.copy()
                         merged_metadata.update(split.metadata) # Add header info
                         # Check length before recursive split
                         if len(split.page_content) > self.chunk_size * 1.1:
                              sub_chunks = self.recursive_splitter.split_text(split.page_content)
                              for sub_chunk in sub_chunks:
                                   split_docs.append(Document(page_content=sub_chunk, metadata=merged_metadata.copy()))
                         else:
                              split_docs.append(Document(page_content=split.page_content, metadata=merged_metadata))

                except Exception as e:
                    logger.warning(f"Header splitting failed for {original_metadata.get('file_name')}, falling back: {e}")
                    # Fallback if header splitting errors
                    split_docs = self.recursive_splitter.split_documents([doc])
            else:
                # Use recursive splitting for PDFs, PPTX, or DOCX without clear headers
                split_docs = self.recursive_splitter.split_documents([doc])

            # Assign chunk IDs relative to the original document
            for i, chunk in enumerate(split_docs):
                 chunk.metadata["chunk_in_doc_id"] = i
                 # Ensure metadata is inherited (recursive_splitter should do this, but double check)
                 for k, v in original_metadata.items():
                      if k not in chunk.metadata:
                           chunk.metadata[k] = v
                 all_chunks.append(chunk)

            doc_index += 1

        # Add a global, unique chunk ID across all documents
        for i, chunk in enumerate(all_chunks):
            chunk.metadata["global_chunk_id"] = str(i)
            chunk_hash = hashlib.md5(chunk.page_content.encode()).hexdigest()[:8]
            chunk.metadata["chunk_hash_prefix"] = chunk_hash # Useful for identification

        logger.info(f"Total chunks created after splitting: {len(all_chunks)}")
        return all_chunks


class QueryProcessor:
    """Advanced query processing for better retrieval performance"""

    def __init__(self, llm):
        self.llm = llm
        # Combine stopwords from both languages
        try:
            self.stop_words = set(stopwords.words('english')).union(set(stopwords.words('french')))
        except Exception as e:
             logger.warning(f"Could not load NLTK stopwords, proceeding without them: {e}")
             self.stop_words = set()

    # Make sync for easier integration first, refactor to async later if needed
    def expand_query(self, query: str, config: Optional[RunnableConfig] = None) -> str:
        """Expand query with related terms using LLM (Synchronous version)"""
        expansion_prompt = PromptTemplate.from_template(
            """Analyze this search query: "{query}"
            Generate 3-5 additional relevant search keywords or synonyms that preserve the original intent but broaden the search possibilities.
            Focus on extracting key concepts and adding variations. Examples: 'cost reduction' -> 'expense saving, budget optimization, efficiency improvement'; 'network security protocols' -> 'firewall configuration, VPN standards, TLS encryption, intrusion detection systems'.
            Return ONLY the additional keywords separated by spaces, nothing else. Do not repeat terms from the original query."""
        )
        logger.debug(f"Attempting query expansion for: {query}")
        try:
            chain = expansion_prompt | self.llm | StrOutputParser()
            expanded_terms_str = chain.invoke({"query": query}, config=config)

            # Clean up and append
            expanded_terms = [term.strip() for term in expanded_terms_str.split() if term.strip() and term.lower() not in self.stop_words]
            # Limit number of added terms
            expanded_terms = expanded_terms[:5]

            if expanded_terms:
                 expanded_query = f"{query} {' '.join(expanded_terms)}"
                 logger.info(f"Query expanded: '{query}' -> '{expanded_query}'")
                 return expanded_query
            else:
                 logger.debug("No useful expansion terms generated.")
                 return query
        except Exception as e:
            logger.warning(f"Query expansion failed: {e}", exc_info=True)
            return query

    # Decomposition is more complex, let's keep it simple for now
    # async def decompose_query(self, query: str) -> List[str]: ...


# --- Enhanced RAG System ---

class EnhancedRAG:
    """Enhanced RAG system with advanced retrieval, processing and generation"""

    def __init__(self,
                 data_dir: Path = DATA_DIR,
                 db_dir: Path = DB_DIR,
                 rebuild: bool = False,
                 embedding_model: str = EMBEDDING_MODEL,
                 llm_model: str = LLM_MODEL,
                 reranker_model: str = RERANKER_MODEL):

        self.data_dir = data_dir
        self.db_dir = db_dir
        self.rebuild = rebuild
        self.embedding_model_name = embedding_model
        self.llm_model_name = llm_model
        self.reranker_model_name = reranker_model

        # Ensure directories exist
        self.data_dir.mkdir(parents=True, exist_ok=True)
        self.db_dir.mkdir(parents=True, exist_ok=True)
        EMBEDDING_CACHE_DIR.mkdir(parents=True, exist_ok=True) # Create embedding cache dir

        logger.info(f"Initializing Enhanced RAG system with:")
        logger.info(f"  - Data Dir: {self.data_dir}")
        logger.info(f"  - DB Dir: {self.db_dir}")
        logger.info(f"  - Rebuild: {self.rebuild}")
        logger.info(f"  - Embedding Model: {self.embedding_model_name}")
        logger.info(f"  - LLM Model: {self.llm_model_name}")
        logger.info(f"  - Reranker Model: {self.reranker_model_name}")

        # Initialize components
        self._initialize_components()

        # Initialize vector store and retriever
        self.vectorstore, self.bm25_retriever, self.all_chunks_for_bm25 = self._load_or_build_vectorstore()

        if self.vectorstore is None:
            raise RuntimeError("Failed to initialize vector store.")

        # Setup hybrid retriever pipeline
        self.full_retriever = self._create_full_retriever_pipeline()

        # Initialize query processor (needs LLM)
        self.query_processor = QueryProcessor(self.llm)

        # Response cache
        self.response_cache = ResponseCache(RESULT_CACHE_FILE)

        # Create RAG chain (needs query processor, retriever, llm)
        self.rag_chain = self._create_rag_chain()
        logger.info("Enhanced RAG system initialization complete.")

    def _initialize_components(self):
        """Initialize core components: LLM, Embeddings, Reranker"""
        logger.info("Initializing core components...")
        try:
            # LLM (Ollama)
            self.llm = OllamaLLM(model=self.llm_model_name, temperature=0.1) # Low temp for factuality

            # Embeddings (HuggingFace with Cache)
            logger.info(f"Setting up embedding cache at: {EMBEDDING_CACHE_DIR}")
            fs_store = LocalFileStore(str(EMBEDDING_CACHE_DIR))
            base_embeddings = HuggingFaceEmbeddings(
                 model_name=self.embedding_model_name,
                 model_kwargs={'device': GPU_DEVICE}
            )
            # Use the Langchain CacheBackedEmbeddings with the store
            self.cached_embedder = LangchainCacheBackedEmbeddings.from_bytes_store(
                underlying_embeddings=base_embeddings,
                document_embedding_cache=fs_store,
                namespace=self.embedding_model_name # Separate cache per model
            )
            logger.info(f"Initialized Embeddings: {self.embedding_model_name} (Cached)")

            # Reranker (CrossEncoder)
            # Load the model using sentence-transformers
            cross_encoder = HuggingFaceCrossEncoder(model_name="BAAI/bge-reranker-base", model_kwargs={'device': 'cpu'})

            # Wrap it for use with LangChain compressor
            
            self.reranker_transformer = CrossEncoderReranker(model=cross_encoder, top_n=RERANKER_TOP_N)
            logger.info(f"Initialized Reranker: {self.reranker_model_name} (top_n={RERANKER_TOP_N})")

            # Document Processor
            self.doc_processor = DocumentProcessor(self.data_dir)

            # Document Splitter
            self.splitter = EnhancedDocumentSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)

        except Exception as e:
            logger.critical(f"Core component initialization failed: {e}", exc_info=True)
            raise

    def _get_current_file_metadata(self) -> Dict[str, float]:
        """Gets modification times for files in data_dir"""
        # (Keep user's implementation - seems okay)
        metadata = {}
        supported_extensions = {".pdf", ".pptx", ".docx"}
        try:
            for file_path in self.data_dir.rglob("*"):
                if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                    try:
                        relative_path = file_path.relative_to(self.data_dir).as_posix()
                        metadata[relative_path] = file_path.stat().st_mtime
                    except Exception as e:
                        logger.warning(f"Could not get metadata for file {file_path}: {e}")
        except Exception as e:
            logger.error(f"Error scanning data directory {self.data_dir}: {e}")
        return metadata

    def _save_chunk_cache(self, chunks: List[Document]):
        """Saves processed chunks to a file for faster BM25 reloading."""
        logger.info(f"Saving {len(chunks)} chunks to cache file: {CHUNK_CACHE_FILE}")
        try:
             # Ensure directory exists
             CHUNK_CACHE_FILE.parent.mkdir(parents=True, exist_ok=True)
             with open(CHUNK_CACHE_FILE, 'w') as f:
                  # Store as list of dicts for JSON compatibility
                  json.dump([doc.dict() for doc in chunks], f)
             logger.info("Chunk cache saved successfully.")
        except Exception as e:
             logger.error(f"Failed to save chunk cache: {e}", exc_info=True)

    def _load_chunk_cache(self) -> Optional[List[Document]]:
         """Loads processed chunks from the cache file."""
         if CHUNK_CACHE_FILE.exists():
              logger.info(f"Loading chunks from cache file: {CHUNK_CACHE_FILE}")
              try:
                   with open(CHUNK_CACHE_FILE, 'r') as f:
                        chunks_data = json.load(f)
                   # Reconstruct Document objects
                   chunks = [Document(**data) for data in chunks_data]
                   logger.info(f"Loaded {len(chunks)} chunks from cache.")
                   return chunks
              except Exception as e:
                   logger.error(f"Failed to load chunk cache: {e}. Will rebuild.", exc_info=True)
                   # Optionally remove corrupted cache file
                   # CHUNK_CACHE_FILE.unlink(missing_ok=True)
                   return None
         else:
              logger.info("Chunk cache file not found.")
              return None


    def _load_or_build_vectorstore(self) -> Tuple[Optional[Chroma], Optional[BM25Retriever], List[Document]]:
        """Loads vectorstore and BM25 if cache is valid, otherwise builds them."""
        current_metadata = self._get_current_file_metadata()
        cache_valid = False
        loaded_chunks_for_bm25 = []

        # --- Cache Validation ---
        if not self.rebuild and CACHE_METADATA_FILE.exists() and CHUNK_CACHE_FILE.exists() and self.db_dir.exists() and any(self.db_dir.iterdir()):
            try:
                with open(CACHE_METADATA_FILE, 'r') as f:
                    cached_metadata = json.load(f)
                if current_metadata == cached_metadata:
                    logger.info("Cache metadata matches. Attempting to load stores and chunks.")
                    cache_valid = True
                else:
                    logger.info("File changes detected. Rebuilding required.")
            except Exception as e:
                logger.warning(f"Cache metadata validation failed: {e}. Rebuilding.")
        elif self.rebuild:
             logger.info("Rebuild explicitly requested.")
        else:
             logger.info("Cache files or DB directory missing. Building required.")

        # --- Cleanup on Rebuild ---
        if self.rebuild or not cache_valid:
             logger.info(f"Clearing existing index directory: {self.db_dir}")
             if self.db_dir.exists(): shutil.rmtree(self.db_dir)
             if CACHE_METADATA_FILE.exists(): CACHE_METADATA_FILE.unlink()
             if CHUNK_CACHE_FILE.exists(): CHUNK_CACHE_FILE.unlink()
             self.db_dir.mkdir(parents=True, exist_ok=True)

        # --- Load from Cache ---
        if cache_valid:
            try:
                logger.info(f"Loading Chroma DB from: {self.db_dir}")
                vectorstore = Chroma(
                    persist_directory=str(self.db_dir),
                    embedding_function=self.cached_embedder # Use cached embedder
                )
                logger.info("Chroma DB loaded successfully.")

                # Load chunks from cache for BM25
                loaded_chunks_for_bm25 = self._load_chunk_cache()
                if not loaded_chunks_for_bm25:
                     logger.warning("Could not load chunks from cache for BM25. BM25 will be unavailable.")
                     return vectorstore, None, [] # Return vectorstore, but no BM25

                bm25_retriever = BM25Retriever.from_documents(loaded_chunks_for_bm25)
                bm25_retriever.k = RETRIEVER_SPARSE_K # Set K for BM25
                logger.info("BM25 retriever initialized from cached chunks.")
                return vectorstore, bm25_retriever, loaded_chunks_for_bm25

            except Exception as e:
                logger.error(f"Error loading from cache/disk: {e}. Rebuilding.", exc_info=True)
                # Fall through to build phase after cleanup
                logger.info(f"Clearing potentially corrupt index directory: {self.db_dir}")
                if self.db_dir.exists(): shutil.rmtree(self.db_dir)
                if CACHE_METADATA_FILE.exists(): CACHE_METADATA_FILE.unlink()
                if CHUNK_CACHE_FILE.exists(): CHUNK_CACHE_FILE.unlink()
                self.db_dir.mkdir(parents=True, exist_ok=True)

        # --- Build Phase ---
        logger.info("Building new vector store and indexes...")
        all_documents = self._process_documents()

        if not all_documents:
            logger.warning(f"No documents processed. Creating empty stores.")
            vectorstore = Chroma.from_documents(
                 documents=[], embedding=self.cached_embedder, persist_directory=str(self.db_dir)
            )
            self._save_cache_metadata(current_metadata) # Save metadata even if empty
            self._save_chunk_cache([]) # Save empty chunk cache
            return vectorstore, None, [] # No BM25 possible

        # Split documents into chunks
        chunks = self.splitter.split_documents(all_documents)
        if not chunks:
            logger.error("No chunks created after splitting. Cannot build stores.")
            # Create empty stores as fallback
            vectorstore = Chroma.from_documents([], embedding=self.cached_embedder, persist_directory=str(self.db_dir))
            self._save_cache_metadata(current_metadata)
            self._save_chunk_cache([])
            return vectorstore, None, []

        logger.info(f"Created {len(chunks)} chunks. Indexing...")

        # Build Chroma DB (using cached embedder automatically handles caching)
        try:
            logger.info("Building Chroma vector store (embeddings will be cached)...")
            start_time = time.time()
            vectorstore = Chroma.from_documents(
                documents=chunks,
                embedding=self.cached_embedder, # Crucial: use the cached embedder
                persist_directory=str(self.db_dir)
            )
            end_time = time.time()
            logger.info(f"Chroma vector store built successfully in {end_time - start_time:.2f} seconds.")
        except Exception as e:
            logger.error(f"Failed to build Chroma DB: {e}", exc_info=True)
            # Cleanup potentially partially created store
            if self.db_dir.exists(): shutil.rmtree(self.db_dir)
            self.db_dir.mkdir(parents=True, exist_ok=True)
            return None, None, []

        # Build BM25 Retriever
        try:
            logger.info("Building BM25 retriever...")
            bm25_retriever = BM25Retriever.from_documents(chunks)
            bm25_retriever.k = RETRIEVER_SPARSE_K # Set K
            logger.info("BM25 retriever built successfully.")
        except Exception as e:
            logger.error(f"Failed to build BM25 retriever: {e}", exc_info=True)
            bm25_retriever = None # Proceed without BM25

        # Save metadata and chunks cache
        self._save_cache_metadata(current_metadata)
        self._save_chunk_cache(chunks) # Save chunks for next load

        return vectorstore, bm25_retriever, chunks


    def _process_documents(self) -> List[Document]:
        """Process all documents in the data directory"""
        # (Keep user's implementation - seems okay)
        all_documents = []
        supported_extensions = {".pdf", ".pptx", ".docx"}
        file_paths_to_process = [
            fp for fp in self.data_dir.rglob("*")
            if fp.is_file() and fp.suffix.lower() in supported_extensions
        ]

        if not file_paths_to_process:
            logger.warning(f"No supported documents found in {self.data_dir}")
            return all_documents

        logger.info(f"Processing {len(file_paths_to_process)} files...")
        # Consider parallel processing for large numbers of files
        # with ThreadPoolExecutor() as executor:
        #     results = executor.map(self._process_single_file, file_paths_to_process)
        #     for result_list in results:
        #         all_documents.extend(result_list)
        for file_path in file_paths_to_process:
             all_documents.extend(self._process_single_file(file_path))

        logger.info(f"Finished processing documents. Total initial docs: {len(all_documents)}")
        return all_documents

    def _process_single_file(self, file_path: Path) -> List[Document]:
         """Helper to process one file based on extension."""
         suffix = file_path.suffix.lower()
         if suffix == ".pdf":
             return self.doc_processor.process_pdf(file_path)
         elif suffix == ".pptx":
             return self.doc_processor.process_pptx(file_path)
         elif suffix == ".docx":
             return self.doc_processor.process_docx(file_path)
         else:
             logger.debug(f"Skipping unsupported file type: {file_path.name}")
             return []

    def _save_cache_metadata(self, metadata: Dict[str, float]):
        """Save file metadata cache"""
        # Ensure directory exists
        CACHE_METADATA_FILE.parent.mkdir(parents=True, exist_ok=True)
        try:
            with open(CACHE_METADATA_FILE, 'w') as f:
                json.dump(metadata, f, indent=2)
            logger.info(f"Metadata cache saved to {CACHE_METADATA_FILE}")
        except Exception as e:
            logger.error(f"Failed to save metadata cache: {e}")

    def _create_full_retriever_pipeline(self) -> Runnable:
        """Creates the full retrieval pipeline: Ensemble -> Rerank -> Compress"""
        logger.info("Creating full retriever pipeline...")

        if self.vectorstore is None:
            raise ValueError("Vectorstore not initialized.")

        # 1. Dense Retriever
        dense_retriever = self.vectorstore.as_retriever(search_kwargs={"k": RETRIEVER_DENSE_K})

        # 2. Base Retriever (Ensemble or Dense Only)
        if self.bm25_retriever is None:
            logger.warning("BM25 retriever unavailable. Using only dense retriever as base.")
            base_retriever = dense_retriever
        else:
            logger.info("Creating Ensemble Retriever (Dense + BM25)")
            # Weights favour dense slightly, tune based on evaluation
            ensemble_weights = [0.6, 0.4]
            self.bm25_retriever.k = RETRIEVER_SPARSE_K # Ensure K is set
            dense_retriever.search_kwargs['k'] = RETRIEVER_DENSE_K # Ensure K is set

            base_retriever = EnsembleRetriever(
                retrievers=[dense_retriever, self.bm25_retriever],
                weights=ensemble_weights,
                search_kwargs={'k': ENSEMBLE_RANK_K} # Return top K after RRF fusion
            )

        # 3. Reranking Retriever (using the CrossEncoder model)
        # Wrap the reranker transformer in ContextualCompressionRetriever
        reranking_retriever = ContextualCompressionRetriever(
            base_compressor=self.reranker_transformer, # Initialized in __init__
            base_retriever=base_retriever # Takes output from ensemble/dense
        )
        logger.info("Reranking step configured.")

        # 4. Contextual Compression Retriever (using LLM)
        # Wrap the LLM extractor in another ContextualCompressionRetriever
        # Note: Using the main LLM for compression can be slow.
        # Consider a smaller/faster LLM instance if latency is an issue.
        final_compression_retriever = ContextualCompressionRetriever(
            base_compressor=LLMChainExtractor.from_llm(self.llm), # Uses main LLM
            base_retriever=reranking_retriever # Takes output from reranker
        )
        logger.info("LLM Contextual Compression step configured.")

        return final_compression_retriever


    def _create_rag_chain(self) -> Runnable:
        """Creates the final RAG chain using LCEL"""
        logger.info("Creating RAG chain using LCEL...")

        # Prompt template for the final answer generation
        # (Using a slightly modified version of the user's template)
        prompt_template = ChatPromptTemplate.from_template(
            """You are an expert AI assistant. Answer the question based **only** on the provided context and conversation history. Be concise and factual. If the information isn't available in the context, clearly state that.

CONVERSATION HISTORY:
{chat_history}

RELEVANT CONTEXT (Filtered and Re-ranked):
{context}

Question: {question}

Answer:"""
        )

        def format_docs(docs: List[Document]) -> str:
            """Formats final list of documents for the prompt"""
            if not docs:
                return "No relevant context found after filtering."
            # Add source markers for clarity in context (optional)
            # return "\n\n".join([f"Source: {d.metadata.get('file_name', 'N/A')} P:{d.metadata.get('page_number', 'N/A')}\n{d.page_content}" for d in docs])
            return "\n\n".join(doc.page_content for doc in docs)

        def format_history(chat_history_tuples: List[Tuple[str,str]]) -> str:
             """Formats chat history tuples into a simple string."""
             if not chat_history_tuples:
                  return "No conversation history."
             return "\n".join([f"Human: {q}\nAssistant: {a}" for q, a in chat_history_tuples])

        # Define the RAG pipeline using LCEL
        rag_chain_pipeline = (
            RunnableParallel( # Step 1: Prepare initial inputs
                {
                    "question": itemgetter("question"),
                    "chat_history_tuples": itemgetter("chat_history"), # Keep original tuples
                    "config": RunnablePassthrough() # Pass config through
                }
            )
            | RunnableParallel( # Step 2: Expand query and format history
                 {
                     # Expand query using QueryProcessor (sync version)
                     "expanded_question": RunnableLambda(lambda x: self.query_processor.expand_query(x["question"], x["config"])),
                     "chat_history": lambda x: format_history(x["chat_history_tuples"]),
                     "original_question": itemgetter("question"), # Pass original q through
                     "config": itemgetter("config")
                 }
             )
            | RunnableParallel( # Step 3: Retrieve documents using expanded query
                {
                    # Retrieve using the full pipeline (ensemble->rerank->compress)
                    "final_docs": RunnableLambda(lambda x: self.full_retriever.invoke(x["expanded_question"], config=x["config"])),
                    "question": itemgetter("original_question"), # Use original question for final prompt
                    "chat_history": itemgetter("chat_history"),
                    "config": itemgetter("config")
                }
            )
            | RunnableParallel( # Step 4: Prepare context and inputs for LLM
                 {
                     "context": lambda x: format_docs(x["final_docs"]),
                     "question": itemgetter("question"),
                     "chat_history": itemgetter("chat_history"),
                     "final_docs": itemgetter("final_docs") # Pass final docs for output
                 }
            )
            | RunnableParallel( # Step 5: Generate answer and pass docs
                 {
                      "answer": prompt_template | self.llm | StrOutputParser(),
                      "docs": itemgetter("final_docs")
                 }
            )
        )

        logger.info("LCEL RAG chain created.")
        return rag_chain_pipeline

    # Make query async to accommodate potential async steps later
    async def query(self, question: str, chat_history: List[Tuple[str, str]] = []) -> Dict[str, Any]:
        """Query the RAG system (asynchronous)"""
        logger.info(f"Received query: '{question[:100]}...'")
        if not hasattr(self, 'rag_chain') or self.rag_chain is None:
            logger.error("RAG chain not initialized.")
            return {"answer": "Error: The RAG system is not ready.", "sources_metadata": []}

        # Check response cache first
        history_hash = self.response_cache.create_history_hash(chat_history)
        cached_response = self.response_cache.get(question, history_hash)
        if cached_response:
            return cached_response # Return cached result directly

        logger.info("Cache miss. Executing RAG chain...")
        start_time = time.time()
        try:
            # Prepare input dictionary for the chain
            chain_input = {"question": question, "chat_history": chat_history}

            # Use ainvoke for the async chain
            # Pass config for potential callbacks/tracing
            config = RunnableConfig(configurable={"callbacks": []}) # Add callbacks if needed
            result = await self.rag_chain.ainvoke(chain_input, config=config)

            answer = result.get("answer", "Error: No answer content generated.")
            final_docs = result.get("docs", [])

            formatted_sources = self._format_source_citation(final_docs)
            full_response = f"{answer}\n\n{formatted_sources}".strip()

            end_time = time.time()
            logger.info(f"RAG chain execution time: {end_time - start_time:.2f} seconds.")
            logger.info(f"Generated answer length: {len(answer)}, Final context docs: {len(final_docs)}")

            response_dict = {
                "answer": full_response,
                "sources_metadata": [doc.metadata for doc in final_docs]
            }

            # Cache the new response
            self.response_cache.set(question, history_hash, response_dict)

            return response_dict

        except Exception as e:
            logger.error(f"Error during async RAG query execution: {e}", exc_info=True)
            return {"answer": f"Sorry, an error occurred during processing: {type(e).__name__}", "sources_metadata": []}

    def _format_source_citation(self, docs: List[Document]) -> str:
        """Format citations based on final document metadata"""
        if not docs:
            return ""

        sources_list = []
        # Keep track of file appearances to potentially group citations
        file_sources = {}

        for doc in docs:
            metadata = doc.metadata
            file_name = metadata.get("file_name", "Unknown Source")
            source_key = f"{file_name}" # Base key on filename

            page_num = metadata.get("page_number")
            slide_num = metadata.get("slide_number")
            chunk_prefix = metadata.get("chunk_hash_prefix", "") # Use hash prefix if available

            location_info = []
            if page_num: location_info.append(f"Pg {page_num}")
            if slide_num: location_info.append(f"Slide {slide_num}")
            # Add chunk info for better localization if needed
            # if chunk_prefix: location_info.append(f"Chunk ~{chunk_prefix}")

            location_str = ", ".join(location_info)
            if location_str:
                 source_key += f" ({location_str})"

            # Store details per file
            if file_name not in file_sources:
                file_sources[file_name] = {"details": set()}
            if location_str:
                file_sources[file_name]["details"].add(location_str)

        # Format the citations
        formatted_list = []
        source_index = 1
        for file_name, data in sorted(file_sources.items()):
            details_str = ""
            if data["details"]:
                 # Sort details for consistent output (e.g., sort page numbers)
                 sorted_details = sorted(list(data["details"]), key=lambda x: int(re.search(r'\d+', x).group()) if re.search(r'\d+', x) else 0)
                 details_str = f" (Locations: {', '.join(sorted_details)})"
            formatted_list.append(f"- Source {source_index}: {file_name}{details_str}")
            source_index += 1


        if not formatted_list:
             return ""

        return "**Sources Consulted (Filtered & Re-ranked):**\n" + "\n".join(formatted_list)


# --- Main Execution & Gradio Interface ---
async def run_query_async(rag_system, message, history):
     """Wrapper to run the async query method"""
     return await rag_system.query(message, history)

def main():
    parser = argparse.ArgumentParser(description="Enhanced RAG System")
    parser.add_argument("--rebuild", action="store_true", help="Force rebuild of indexes")
    parser.add_argument("--data-dir", type=str, default=str(DATA_DIR), help="Source documents directory")
    parser.add_argument("--db-dir", type=str, default=str(DB_DIR), help="Vector database directory")
    # Add args for models if needed
    # parser.add_argument("--llm-model", type=str, default=LLM_MODEL)
    # parser.add_argument("--embedding-model", type=str, default=EMBEDDING_MODEL)
    # parser.add_argument("--reranker-model", type=str, default=RERANKER_MODEL)

    args = parser.parse_args()
    data_path = Path(args.data_dir).resolve()
    db_path = Path(args.db_dir).resolve()

    if not data_path.is_dir():
        print(f"Error: Data directory not found at '{data_path}'")
        sys.exit(1)

    try:
        logger.info(f"Initializing Enhanced RAG: Data='{data_path}', DB='{db_path}', Rebuild={args.rebuild}")
        # Pass model args if defined via CLI
        enhanced_rag = EnhancedRAG(
            data_dir=data_path,
            db_dir=db_path,
            rebuild=args.rebuild,
            # llm_model=args.llm_model, # etc.
        )
    except Exception as e:
        logger.critical(f"Initialization failed: {e}", exc_info=True)
        sys.exit(1)

    # --- Gradio Interface ---
    with gr.Blocks(theme=ThemeBase(), title="Enhanced RAG Assistant") as demo:
        gr.Markdown("# Enhanced RAG Assistant")
        gr.Markdown(f"Querying documents from `{data_path}` using advanced retrieval.")

        chatbot_ui = gr.Chatbot(label="Conversation", height=600, show_copy_button=True, bubble_full_width=False)
        msg_input = gr.Textbox(label="Your question:", placeholder="Ask about the documents...", lines=3, scale=7) # Larger input box

        with gr.Row():
             submit_button = gr.Button("Send", variant="primary", scale=1)
             clear_button = gr.Button("Clear Conversation", variant="secondary", scale=1)

        # Gradio state to store chat history List[Tuple[str, str]]
        chat_state = gr.State([])

        async def handle_submit(user_message: str, history: List[Tuple[str, str]]):
             """Handles user submission, calls async query, updates history"""
             # Append user message immediately for responsiveness
             history.append([user_message, None])
             yield history # Update UI to show user message

             # Call the async query function
             response_dict = await enhanced_rag.query(user_message, history[:-1]) # Pass history *before* current turn
             bot_message = response_dict["answer"]

             # Update the history with the bot's response
             history[-1][1] = bot_message
             yield history # Update UI with bot response

        # Connect Gradio events to the async handler
        submit_button.click(
             handle_submit,
             inputs=[msg_input, chat_state],
             outputs=[chat_state],
             queue=True # Enable queuing for async handling
        ).then(
            lambda: gr.update(value=""), None, [msg_input], queue=False # Clear input box after submit
        )

        msg_input.submit(
             handle_submit,
             inputs=[msg_input, chat_state],
             outputs=[chat_state],
             queue=True
        ).then(
            lambda: gr.update(value=""), None, [msg_input], queue=False # Clear input box after submit
        )

        # Update chatbot UI based on chat_state changes
        chat_state.change(lambda state: state, inputs=[chat_state], outputs=[chatbot_ui], queue=True)

        # Clear button functionality
        def clear_chat():
            return [], [] # Return empty lists for both chatbot UI and chat_state

        clear_button.click(clear_chat, [], [chatbot_ui, chat_state], queue=False)

    logger.info("Launching Gradio interface...")
    # Add server_name="0.0.0.0" to listen on all interfaces if needed
    demo.queue().launch(share=False, server_name="0.0.0.0")


if __name__ == "__main__":
    # Ensure asyncio event loop is running correctly for Gradio async
    # This is often handled automatically by modern Gradio/uvicorn, but explicit setup can be safer
    # if sys.platform == "win32":
    #      asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())
    main()