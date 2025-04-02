__import__('pysqlite3')
import sys
import os
sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')

import os
import argparse
import re
import json
import datetime
import time
import hashlib
import logging
from typing import List, Dict, Optional, Generator, Tuple, Any, AsyncGenerator
from pathlib import Path
from io import BytesIO
import tempfile
import asyncio
from concurrent.futures import ThreadPoolExecutor
from enum import Enum

from langchain_chroma import Chroma
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.documents import Document
from langchain_core.runnables import RunnablePassthrough
from langchain_ollama import OllamaLLM
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from pptx import Presentation
from docx import Document as WordDocument
from PIL import Image
import torch
from transformers import AutoProcessor, AutoModelForVision2Seq

import gradio as gr
from gradio.themes.utils import colors, sizes
from gradio.themes import Base as ThemeBase

# --- Configuration ---
EMBEDDING_MODEL = "sentence-transformers/all-mpnet-base-v2"  # Keep reliable embeddings
DOC_PROCESSING_MODEL = "ds4sd/SmolDocling-256M-preview"     # For document processing
LLM_MODEL = "gemma3:27b-it-q8_0" 
DB_DIR = Path("vectore/chroma_rag_db")
DATA_DIR = Path("GO 2024")
CACHE_FILE = DB_DIR / "rag_cache_metadata.json"
RETRIEVER_K = 7
RETRIEVER_FETCH_K = 15

# --- Document Element Types ---
class DocElementType(Enum):
    TEXT = "text"
    TABLE = "table"
    FIGURE = "figure"
    CODE = "code"
    FORMULA = "formula"
    HEADER = "header"

# --- Logging Setup ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- SmolDocling Processor Class ---
class SmolDoclingProcessor:
    def __init__(self, model_name: str = DOC_PROCESSING_MODEL):
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        logging.info(f"Loading SmolDocling processor on {self.device}...")
        
        try:
            # Initialize processor and model
            self.processor = AutoProcessor.from_pretrained(model_name)
            self.model = AutoModelForVision2Seq.from_pretrained(model_name).to(self.device)
            logging.info("SmolDocling processor loaded successfully")
        except Exception as e:
            logging.error(f"Failed to load SmolDocling: {e}")
            raise

    def process_document_page(self, image_path: Path) -> List[Dict]:
        """Process a single document page into structured elements"""
        try:
            # Load and prepare image
            image = Image.open(image_path).convert("RGB")
            
            # Process with SmolDocling
            inputs = self.processor(images=image, return_tensors="pt").to(self.device)
            
            generated_ids = self.model.generate(
                **inputs,
                max_new_tokens=8192,
                bad_words_ids=[[self.processor.tokenizer.unk_token_id]],
                output_scores=True,
                return_dict_in_generate=True
            )
            
            # Decode and parse the output
            generated_text = self.processor.batch_decode(
                generated_ids.sequences, 
                skip_special_tokens=True
            )[0]
            
            return self._parse_doctags(generated_text)
        except Exception as e:
            logging.error(f"Error processing document page: {e}")
            return []

    def _parse_doctags(self, doctags_text: str) -> List[Dict]:
        """Convert raw DocTags output to structured elements"""
        elements = []
        current_element = None
        
        for line in doctags_text.split('\n'):
            if line.startswith('<') and line.endswith('>'):
                # Close previous element
                if current_element:
                    elements.append(current_element)
                
                # Start new element
                tag_parts = line[1:-1].split()
                if len(tag_parts) < 2:
                    continue
                    
                element_type = tag_parts[0].lower()
                coordinates = list(map(float, tag_parts[1:5]))
                
                current_element = {
                    "type": element_type,
                    "bbox": coordinates,
                    "content": "",
                    "metadata": {}
                }
                
                # Parse additional metadata if present
                if len(tag_parts) > 5:
                    for part in tag_parts[5:]:
                        if '=' in part:
                            key, val = part.split('=', 1)
                            current_element["metadata"][key] = val
            elif current_element:
                current_element["content"] += line + '\n'
        
        if current_element:
            elements.append(current_element)
            
        return elements

# --- Document RAG Class ---
class DocumentRAG:
    def __init__(self, data_dir: Path = DATA_DIR, db_dir: Path = DB_DIR, rebuild: bool = True):
        self.data_dir = data_dir
        self.db_dir = db_dir

        # Ensure parent directories exist
        self.data_dir.mkdir(parents=True, exist_ok=True)
        self.db_dir.mkdir(parents=True, exist_ok=True)

        # Initialize models
        try:
            logging.info("Initializing embedding model...")
            self.embeddings = HuggingFaceEmbeddings(
                model_name=EMBEDDING_MODEL,
                model_kwargs={'device': 'cuda' if torch.cuda.is_available() else 'cpu'}
            )
            
            logging.info("Initializing LLM...")
            self.llm = OllamaLLM(model=LLM_MODEL)
            
            logging.info("Initializing SmolDocling processor...")
            self.smol_docling = SmolDoclingProcessor()
        except Exception as e:
            logging.error(f"Model initialization failed: {e}")
            raise

        if rebuild and self.db_dir.exists():
            logging.info(f"Rebuild requested. Removing existing vector store: {self.db_dir}")
            import shutil
            shutil.rmtree(self.db_dir)
            if CACHE_FILE.exists():
                CACHE_FILE.unlink()
            self.db_dir.mkdir(parents=True, exist_ok=True)

        self.vectorstore = self._load_or_build_vectorstore()

        self.retriever = self.vectorstore.as_retriever(
            search_type="mmr",
            search_kwargs={"k": RETRIEVER_K, "fetch_k": RETRIEVER_FETCH_K}
        )

        self.rag_chain = self._create_rag_chain()

    def _get_current_file_metadata(self) -> Dict[str, float]:
        """Gets modification times for all supported files in data_dir."""
        metadata = {}
        supported_extensions = {".pdf", ".pptx", ".docx"}
        try:
            for file_path in self.data_dir.rglob("*"):
                if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                    try:
                        relative_path = file_path.relative_to(self.data_dir).as_posix()
                        metadata[relative_path] = file_path.stat().st_mtime
                    except Exception as e:
                         logging.warning(f"Could not get metadata for file {file_path}: {e}")
        except Exception as e:
            logging.error(f"Error scanning data directory {self.data_dir}: {e}")
        return metadata

    def _load_or_build_vectorstore(self) -> Chroma:
        """Loads vectorstore if cache is valid, otherwise builds it."""
        current_metadata = self._get_current_file_metadata()
        cache_valid = False

        if CACHE_FILE.exists() and self.db_dir.exists() and any(self.db_dir.iterdir()):
            try:
                with open(CACHE_FILE, 'r') as f:
                    cached_metadata = json.load(f)
                if current_metadata == cached_metadata:
                    logging.info("Cache metadata matches current files. Loading existing vector store.")
                    cache_valid = True
                else:
                    logging.info("File changes detected (add/remove/modify). Rebuilding vector store.")
                    logging.info(f"Removing outdated vector store: {self.db_dir}")
                    import shutil
                    shutil.rmtree(self.db_dir)
                    self.db_dir.mkdir(parents=True, exist_ok=True)

            except Exception as e:
                logging.warning(f"Could not validate cache file {CACHE_FILE}. Rebuilding vector store. Error: {e}")
                if self.db_dir.exists():
                     logging.info(f"Removing potentially corrupt vector store: {self.db_dir}")
                     import shutil
                     shutil.rmtree(self.db_dir)
                     self.db_dir.mkdir(parents=True, exist_ok=True)

        if cache_valid:
            try:
                return Chroma(
                    persist_directory=str(self.db_dir),
                    embedding_function=self.embeddings
                )
            except Exception as e:
                logging.error(f"Error loading vector store from {self.db_dir}, rebuilding. Error: {e}")

        logging.info("Building new vector store...")
        all_documents = []
        supported_extensions = {".pdf", ".pptx", ".docx"}
        file_paths_to_process = [
            fp for fp in self.data_dir.rglob("*")
            if fp.is_file() and fp.suffix.lower() in supported_extensions
        ]

        if not file_paths_to_process:
             logging.warning(f"No supported documents found in {self.data_dir}. Vector store will be empty.")
             vectorstore = Chroma(
                 embedding_function=self.embeddings,
                 persist_directory=str(self.db_dir)
             )
             current_metadata = {}
             with open(CACHE_FILE, 'w') as f:
                json.dump(current_metadata, f)
             return vectorstore

        for file_path in file_paths_to_process:
            if file_path.suffix.lower() == ".pdf":
                all_documents.extend(self._process_pdf_files(file_path))
            elif file_path.suffix.lower() == ".pptx":
                all_documents.extend(self._process_pptx_files(file_path))
            elif file_path.suffix.lower() == ".docx":
                all_documents.extend(self._process_word_files(file_path))

        if not all_documents:
            raise ValueError(f"No content could be extracted from documents in {self.data_dir}. Check file contents and permissions.")

        logging.info(f"Extracted content from {len(file_paths_to_process)} files, resulting in {len(all_documents)} initial document pieces.")

        chunks = self._create_chunks_parallel(all_documents)

        if not chunks:
            raise ValueError("No chunks created after splitting. Check text splitter settings and document content.")

        logging.info(f"Created {len(chunks)} chunks for vector store.")

        vectorstore = Chroma.from_documents(
            documents=chunks,
            embedding=self.embeddings,
            persist_directory=str(self.db_dir)
        )
        logging.info(f"Vector store built successfully with {len(chunks)} chunks and persisted to {self.db_dir}")

        try:
             with open(CACHE_FILE, 'w') as f:
                json.dump(current_metadata, f)
             logging.info(f"Cache metadata saved to {CACHE_FILE}")
        except Exception as e:
             logging.error(f"Failed to save cache metadata to {CACHE_FILE}: {e}")

        return vectorstore

    # --- Enhanced File Processing Methods ---
    def _process_pdf_files(self, file_path: Path) -> List[Document]:
        documents = []
        try:
            logging.info(f"Processing PDF with SmolDocling: {file_path.name}")
            
            # Create temp directory for page images
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                
                # Convert PDF to images (one per page)
                pdf = PdfReader(file_path)
                pdf_metadata_base = {
                    "source": str(file_path.relative_to(self.data_dir)),
                    "file_name": file_path.name,
                    "file_type": ".pdf",
                    "total_pages": len(pdf.pages)
                }
                
                for page_idx, page in enumerate(pdf.pages):
                    page_num = page_idx + 1
                    image_path = temp_path / f"page_{page_num}.png"
                    
                    # Convert page to image
                    try:
                        page_image = page.to_image(resolution=300)
                        page_image.save(image_path)
                    except Exception as e:
                        logging.warning(f"Could not convert page {page_num} to image: {e}")
                        continue
                    
                    # Process with SmolDocling
                    elements = self.smol_docling.process_document_page(image_path)
                    if not elements:
                        logging.debug(f"  No elements extracted from page {page_num}")
                        continue
                    
                    # Create documents for each element type
                    for element in elements:
                        try:
                            element_type = element["type"]
                            content = element["content"].strip()
                            if not content:
                                continue
                                
                            # Enhanced metadata
                            page_metadata = {
                                **pdf_metadata_base,
                                "page_number": page_num,
                                "element_type": element_type,
                                "element_bbox": element["bbox"],
                                **element.get("metadata", {})
                            }
                            
                            # Special handling for different element types
                            if element_type == "table":
                                content = f"TABLE:\n{content}"
                            elif element_type == "figure":
                                content = f"FIGURE: {element.get('metadata', {}).get('caption', '')}\n{content}"
                            
                            doc = Document(
                                page_content=content,
                                metadata=page_metadata
                            )
                            documents.append(doc)
                        except Exception as e:
                            logging.warning(f"Error processing element on page {page_num}: {e}")
                            
        except Exception as e:
            logging.error(f"Error processing PDF {file_path}: {e}", exc_info=True)
            # Fallback to traditional PDF processing if SmolDocling fails
            try:
                logging.info("Attempting fallback PDF processing...")
                pdf = PdfReader(file_path)
                pdf_metadata_base = {
                    "source": str(file_path.relative_to(self.data_dir)),
                    "file_name": file_path.name,
                    "file_type": ".pdf",
                    "total_pages": len(pdf.pages)
                }
                for page_idx, page in enumerate(pdf.pages):
                    page_num = page_idx + 1
                    text = page.extract_text()
                    if not text or not text.strip():
                        logging.debug(f"  Skipping empty page {page_num} in {file_path.name}")
                        continue
                    page_metadata = {**pdf_metadata_base, "page_number": page_num}
                    doc = Document(page_content=text, metadata=page_metadata)
                    documents.append(doc)
            except Exception as fallback_e:
                logging.error(f"Fallback PDF processing also failed for {file_path}: {fallback_e}")
        return documents

    def _process_pptx_files(self, file_path: Path) -> List[Document]:
        documents = []
        try:
            logging.info(f"Processing presentation: {file_path.name}")
            presentation = Presentation(file_path)
            pptx_metadata_base = {
                "source": str(file_path.relative_to(self.data_dir)),
                "file_name": file_path.name,
                "file_type": ".pptx",
                "total_slides": len(presentation.slides)
            }
            for slide_idx, slide in enumerate(presentation.slides):
                slide_num = slide_idx + 1
                all_text = self._extract_all_text_from_slide(slide)
                if not all_text or not all_text.strip():
                    logging.debug(f"  Skipping empty slide {slide_num} in {file_path.name}")
                    continue
                slide_metadata = {**pptx_metadata_base, "slide_number": slide_num}
                doc = Document(page_content=all_text, metadata=slide_metadata)
                documents.append(doc)
        except Exception as e:
            logging.error(f"Error processing PPTX {file_path}: {e}", exc_info=True)
        return documents

    def _process_word_files(self, file_path: Path) -> List[Document]:
        documents = []
        try:
            logging.info(f"Processing Word document: {file_path.name}")
            doc = WordDocument(file_path)
            word_metadata = {
                "source": str(file_path.relative_to(self.data_dir)),
                "file_name": file_path.name,
                "file_type": ".docx",
            }
            full_text = []
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    full_text.append(para.text.strip())
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    cell_texts = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                    if cell_texts:
                        table_rows.append(" | ".join(cell_texts))
                if table_rows:
                    full_text.append("Table:\n" + "\n".join(table_rows))
            combined_text = "\n\n".join(full_text).strip()
            if not combined_text:
                logging.debug(f"  Skipping empty document: {file_path.name}")
                return documents
            doc_entry = Document(page_content=combined_text, metadata=word_metadata)
            documents.append(doc_entry)
            logging.debug(f"  Processed Word doc {file_path.name} ({len(combined_text)} chars)")
        except Exception as e:
            logging.error(f"Error processing DOCX {file_path}: {e}", exc_info=True)
        return documents

    def _extract_all_text_from_slide(self, slide) -> str:
        text_parts = []
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        text_parts.append(text)
                elif shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        cell_texts = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame and cell.text_frame.text.strip()]
                        if cell_texts:
                            table_rows.append(" | ".join(cell_texts))
                    if table_rows:
                        text_parts.append("Table:\n" + "\n".join(table_rows))
        except Exception as e:
            logging.warning(f"Error extracting text from a shape/table on a slide: {e}")
        return "\n\n".join(text_parts).strip()

    def _create_chunks_parallel(self, documents: List[Document]) -> List[Document]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=500,
            chunk_overlap=100,
            separators=["\n\n", "\n", ". ", ", ", " ", ""],
            length_function=len,
        )
        all_chunks = []
        logging.info("Splitting documents into chunks...")
        try:
            all_chunks = splitter.split_documents(documents)
            logging.info(f"Finished splitting into {len(all_chunks)} chunks.")
        except Exception as e:
             logging.error(f"Error during document splitting: {e}", exc_info=True)
             raise
        return all_chunks

    # --- Formatting and Chain Creation ---
    def _format_chat_history(self, chat_history: List[Tuple[str, str]]) -> str:
        if not chat_history:
            return "No conversation history yet."
        formatted_history = "\n".join([f"Human: {q}\nAssistant: {a}" for q, a in chat_history])
        return formatted_history

    def _format_docs(self, docs: List[Document]) -> str:
        if not docs:
            return "No relevant documents found."
        formatted_docs = []
        for i, doc in enumerate(docs):
            metadata = doc.metadata
            file_name = metadata.get("file_name", "Unknown Source")
            content = doc.page_content
            source_info = f"Source {i+1}: {file_name}"
            
            # Enhanced source info with element types
            if metadata.get("file_type") == ".pdf":
                page_num = metadata.get("page_number", "?")
                source_info += f" (Page {page_num}"
                
                element_type = metadata.get("element_type")
                if element_type:
                    source_info += f", {element_type.upper()}"
                
                source_info += ")"
                
            elif metadata.get("file_type") == ".pptx":
                slide_num = metadata.get("slide_number", "?")
                source_info += f" (Slide {slide_num})"
                
            formatted_docs.append(f"{source_info}\nContent: {content}\n---")
        return "\n".join(formatted_docs)

    def _create_rag_chain(self):
        context_prompt = ChatPromptTemplate.from_template(
            """
            Vous êtes un assistant IA chargé de répondre aux questions en vous basant **uniquement** sur le contexte fourni et l'historique de la conversation.

            HISTORIQUE DE LA CONVERSATION:
            {chat_history}

            CONTEXTE RÉCUPÉRÉ (Documents Pertinents):
            {context}

            Directives Strictes:
            1.  Basez votre réponse **exclusivement** sur le CONTEXTE RÉCUPÉRÉ et l'HISTORIQUE DE LA CONVERSATION.
            2.  **Ne supposez rien et n'inventez aucune information.** Si la réponse n'est pas dans le contexte, déclarez explicitement que l'information n'est pas disponible dans les documents fournis.
            3.  Répondez directement à la question posée.
            4.  **Citez vos sources** en utilisant le format `[Source X]` où X est le numéro de la source indiqué dans le CONTEXTE RÉCUPÉRÉ. Intégrez les citations de manière fluide dans votre réponse. Exemple: "Le projet Alpha a débuté en 2023 [Source 1]."
            5.  Si plusieurs sources confirment un point, vous pouvez les citer ensemble, e.g., `[Source 1, Source 3]`.
            6.  Gardez un ton factuel et professionnel.
            7.  Si la question fait référence à la conversation précédente, tenez-en compte.

            Question: {question}

            Réponse (basée uniquement sur le contexte et l'historique):
            """
        )
        rag_chain_with_sources = RunnablePassthrough.assign(
            source_documents=lambda x: self.retriever.invoke(x["question"]),
            question=lambda x: x["question"],
            chat_history=lambda x: x.get("chat_history", [])
        ) | RunnablePassthrough.assign(
            context=lambda x: self._format_docs(x["source_documents"]),
            chat_history_str=lambda x: self._format_chat_history(x["chat_history"])
        ) | {
            "source_documents": lambda x: x["source_documents"],
            "answer": context_prompt | self.llm | StrOutputParser()
        }
        return rag_chain_with_sources

    def _format_source_citation(self, docs: List[Document]) -> str:
        if not docs:
            return ""
        sources_seen = set()
        formatted_sources = []
        for i, doc in enumerate(docs):
            metadata = doc.metadata
            file_name = metadata.get("file_name", "Unknown Source")
            source_key = f"Source {i+1}: {file_name}"
            if metadata.get("file_type") == ".pdf":
                page_num = metadata.get("page_number", "?")
                source_key += f" (Page {page_num}"
                element_type = metadata.get("element_type")
                if element_type:
                    source_key += f", {element_type.upper()}"
            elif metadata.get("file_type") == ".pptx":
                slide_num = metadata.get("slide_number", "?")
                source_key += f" (Slide {slide_num})"
            unique_id = f"{source_key}_{metadata.get('chunk_id', id(doc))}" # Use id as fallback uniqueifier
            if unique_id not in sources_seen:
                formatted_sources.append(source_key)
                sources_seen.add(unique_id)
        if not formatted_sources:
             return ""
        return "\n\n**Sources Consultées:**\n" + "\n".join(f"- {s}" for s in formatted_sources)

    async def query_stream(self, question: str, chat_history: List[Tuple[str, str]]) -> AsyncGenerator[Dict[str, Any], None]:
        """
        Processes a query with chat history and streams the response dictionary asynchronously.
        Yields dictionaries like: {"type": "chunk", "content": "response part..."}
                                  {"type": "sources", "content": "Formatted sources..."}
                                  {"type": "error", "content": "Error message..."}
        """
        try:
            chain_input = {"question": question, "chat_history": chat_history}
            full_response_answer = ""
            streamed_sources = None

            logging.info(f"Invoking RAG chain stream for question: {question[:50]}...")
            async for chunk in self.rag_chain.astream(chain_input):
                 if "answer" in chunk:
                    full_response_answer += chunk["answer"]
                    yield {"type": "chunk", "content": full_response_answer}
                 if "source_documents" in chunk and chunk["source_documents"] and streamed_sources is None:
                     streamed_sources = chunk["source_documents"]

            logging.info("RAG chain stream finished.")

            if streamed_sources is None:
                 logging.warning("Source documents not found or empty in stream.")

            formatted_sources = self._format_source_citation(streamed_sources or [])
            if formatted_sources:
                yield {"type": "sources", "content": formatted_sources}

        except ConnectionError as ce:
            logging.error(f"Connection Error during RAG query: {ce}", exc_info=True)
            yield {"type": "error", "content": f"Erreur de Connexion: Impossible de joindre le modèle LLM. Vérifiez si Ollama est lancé. ({ce})"}
        except TimeoutError as te:
            logging.error(f"Timeout Error during RAG query: {te}", exc_info=True)
            yield {"type": "error", "content": f"Erreur de Timeout: La requête a pris trop de temps. ({te})"}
        except Exception as e:
            logging.error(f"Unexpected error during RAG async query: {e}", exc_info=True)
            yield {"type": "error", "content": f"Erreur Inattendue: {e}"}

    def query(self, question: str, chat_history: List[Tuple[str, str]]) -> Dict[str, Any]:
        try:
            chain_input = {"question": question, "chat_history": chat_history}
            result = self.rag_chain.invoke(chain_input)
            answer = result.get("answer", "Erreur: Aucune réponse générée.")
            source_docs = result.get("source_documents", [])
            formatted_sources = self._format_source_citation(source_docs)
            return {"answer": answer, "sources": formatted_sources}
        except Exception as e:
            logging.error(f"Error during RAG query (non-streaming): {e}", exc_info=True)
            return {"answer": f"Erreur: {e}", "sources": ""}

# --- Gradio Theme ---
class CustomTheme(ThemeBase):
    def __init__(self):
        super().__init__(
            primary_hue=colors.blue,
            secondary_hue=colors.cyan,
            neutral_hue=colors.gray,
            spacing_size=sizes.spacing_md,
            radius_size=sizes.radius_md,
            text_size=sizes.text_md,
        )
    def set_styles(self):
        super().set_styles()
        self.styles.update({
            "button": {"padding": f"{sizes.spacing_sm} {sizes.spacing_md}","border_radius": sizes.radius_md,},
            "button_primary": {"background": f"linear-gradient(to right, {colors.blue[600]}, {colors.blue[500]})","color": colors.white,"_hover": {"background": f"linear-gradient(to right, {colors.blue[700]}, {colors.blue[600]})",}},
            "chatbot": {"border_radius": sizes.radius_lg,"box_shadow": f"0 2px 8px {colors.gray[200]}",},
            "input": {"border_radius": sizes.radius_lg,"box_shadow": f"0 2px 8px {colors.gray[200]}",}
        })

# --- Main Application ---
def main():
    parser = argparse.ArgumentParser(description="Document RAG System with SmolDocling Integration")
    parser.add_argument("--rebuild", action="store_true", help="Force rebuild of the vector database")
    parser.add_argument("--query", type=str, help="Run a specific query in non-interactive mode (outputs JSON)")
    parser.add_argument("--sharepoint", action="store_true", help="Enable SharePoint integration (requires env vars)")
    parser.add_argument("--data-dir", type=str, default=str(DATA_DIR), help=f"Directory containing documents (default: {DATA_DIR})")
    parser.add_argument("--db-dir", type=str, default=str(DB_DIR), help=f"Directory to store vector database (default: {DB_DIR})")
    args = parser.parse_args()

    data_path = Path(args.data_dir)
    db_path = Path(args.db_dir)

    try:
        logging.info(f"Initializing RAG system: Data='{data_path}', DB='{db_path}', Rebuild={args.rebuild}")
        rag = DocumentRAG(
            data_dir=data_path,
            db_dir=db_path,
            rebuild=args.rebuild,
        )
        logging.info("RAG system initialized successfully.")
    except Exception as e:
        logging.error(f"Failed to initialize RAG system: {e}", exc_info=True)
        return

    if args.query:
        logging.info(f"Running single query: {args.query}")
        result = rag.query(args.query, chat_history=[])
        print(json.dumps(result, indent=2, ensure_ascii=False))
        return

    logging.info("Starting Gradio interface...")
    custom_theme = CustomTheme()

    with gr.Blocks(theme=custom_theme, title="Assistant Documentaire RAG") as demo:
        gr.Markdown("""
        # Assistant Documentaire Intelligent
        Posez vos questions sur vos documents. L'assistant répondra en se basant sur les informations trouvées et citera ses sources.
        """)

        chatbot = gr.Chatbot(label="Conversation", height=600, bubble_full_width=False, show_label=False)
        chat_history_state = gr.State([])

        with gr.Row():
            msg = gr.Textbox(label="Votre Question", placeholder="Posez votre question ici...", scale=7, container=False, autofocus=True, show_label=False)
            submit_btn = gr.Button("Envoyer", variant="primary", scale=1)
            clear_btn = gr.Button("Effacer", variant="secondary", scale=1)

        status = gr.Textbox(value="Prêt", label="Statut", interactive=False, max_lines=1)

        def user(message, history_list):
            history_list.append((message, None))
            return "", history_list

        async def bot(history_list):
            if not history_list or history_list[-1][0] is None:
                 yield history_list, gr.update(value="Erreur: Aucune question fournie.")
                 return

            question = history_list[-1][0]
            context_history = history_list[:-1]

            yield history_list, gr.update(value="Recherche et génération en cours...")

            full_answer = ""
            sources_text = ""
            error_text = ""

            try:
                 async for update in rag.query_stream(question, context_history):
                     if update["type"] == "chunk":
                         full_answer = update["content"]
                         history_list[-1] = (question, full_answer + "▌")
                         yield history_list, gr.update(value="Génération...")
                     elif update["type"] == "sources":
                         sources_text = update["content"]
                         history_list[-1] = (question, full_answer + sources_text)
                         yield history_list, gr.update(value="Finalisation...")
                     elif update["type"] == "error":
                         error_text = update["content"]
                         history_list[-1] = (question, error_text)
                         yield history_list, gr.update(value="Erreur")
                         return

                 if sources_text and history_list[-1][1].endswith("▌"):
                      history_list[-1] = (question, history_list[-1][1][:-1])

                 elif not sources_text and history_list[-1][1].endswith("▌"):
                      history_list[-1] = (question, full_answer)

                 if not error_text:
                     yield history_list, gr.update(value="Prêt")

            except Exception as e:
                 logging.error(f"Error in Gradio bot async function: {e}", exc_info=True)
                 error_msg = f"Erreur inattendue dans l'interface: {e}"
                 current_answer = history_list[-1][1]
                 if isinstance(current_answer, str) and current_answer.endswith("▌"):
                      history_list[-1] = (question, current_answer[:-1] + f"\n\n{error_msg}")
                 else:
                      history_list[-1] = (question, error_msg)
                 yield history_list, gr.update(value="Erreur Critique")


        msg.submit(user, [msg, chat_history_state], [msg, chat_history_state], queue=False).then(
            bot, [chat_history_state], [chatbot, status]
        )
        submit_btn.click(user, [msg, chat_history_state], [msg, chat_history_state], queue=False).then(
            bot, [chat_history_state], [chatbot, status]
        )
        clear_btn.click(lambda: ([], []), None, [chatbot, chat_history_state], queue=False).then(
            lambda: "Prêt", None, status
        )

    demo.queue().launch()

if __name__ == "__main__":
    main()